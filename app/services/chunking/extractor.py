"""
Text extraction service.

Extracts raw text from various file formats:
- Plain text (.txt)
- PDF (.pdf)
- Word Document (.docx)
- Markdown (.md)
- HTML (.html)
- CSV (.csv)
- JSON (.json)
- PowerPoint (.pptx)

Each format is handled by a dedicated extractor method.
Unsupported formats raise UnsupportedFileFormatError.
"""

from __future__ import annotations

import csv
import json
import logging
import re
from io import StringIO
from pathlib import Path
from typing import Any

from app.core.exceptions import UnsupportedFileFormatError

logger = logging.getLogger(__name__)

# Mapping of file extensions to MIME types
SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".html": "text/html",
    ".htm": "text/html",
    ".csv": "text/csv",
    ".json": "application/json",
    ".xml": "application/xml",
}

# Extensions that require optional dependencies
OPTIONAL_EXTENSIONS: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class TextExtractorService:
    """Extracts raw text content from files of supported formats.

    Usage:
        extractor = TextExtractorService()
        text = await extractor.extract("path/to/file.pdf")
        # or
        text = await extractor.extract_bytes(file_bytes, extension=".pdf")
    """

    # ── Public API ───────────────────────────────────────────────────

    async def extract(self, file_path: str | Path) -> str:
        """Extract text from a file.

        Args:
            file_path: Path to the file.

        Returns:
            Extracted text content.

        Raises:
            UnsupportedFileFormatError: If the file format is not supported.
        """
        path = Path(file_path)
        extension = path.suffix.lower()

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        content_bytes = await self._read_file(path)
        return await self.extract_bytes(content_bytes, extension)

    async def extract_bytes(self, content: bytes, extension: str) -> str:
        """Extract text from raw bytes with a given file extension.

        Args:
            content: Raw file bytes.
            extension: File extension (e.g., '.pdf', '.txt').

        Returns:
            Extracted text content.

        Raises:
            UnsupportedFileFormatError: If the format is not supported.
        """
        extension = extension.lower()

        if extension in SUPPORTED_EXTENSIONS:
            return await self._extract_builtin(content, extension)
        elif extension in OPTIONAL_EXTENSIONS:
            return await self._extract_optional(content, extension)
        else:
            raise UnsupportedFileFormatError(extension)

    # ── Built-in format extractors ───────────────────────────────────

    async def _extract_builtin(self, content: bytes, extension: str) -> str:
        """Extract text from formats supported without extra dependencies.

        Args:
            content: Raw file bytes.
            extension: File extension.

        Returns:
            Extracted text.
        """
        try:
            if extension == ".txt":
                return self._extract_plain_text(content)
            elif extension == ".md":
                return self._extract_markdown(content)
            elif extension in (".html", ".htm"):
                return self._extract_html(content)
            elif extension == ".csv":
                return self._extract_csv(content)
            elif extension == ".json":
                return self._extract_json(content)
            elif extension == ".xml":
                return self._extract_xml(content)
            else:
                # Fallback: return raw text
                text = content.decode("utf-8", errors="replace")
                return text.strip()
        except Exception as e:
            raise ValueError(f"Failed to extract text from {extension} file: {e}") from e

    @staticmethod
    def _extract_plain_text(content: bytes) -> str:
        """Extract text from a plain text file."""
        return content.decode("utf-8", errors="replace").strip()

    @staticmethod
    def _extract_markdown(content: bytes) -> str:
        """Extract text from a Markdown file.

        Strips Markdown formatting to get clean text.
        """
        text = content.decode("utf-8", errors="replace")
        # Remove code blocks
        text = re.sub(r"```[\s\S]*?```", "", text)
        # Remove images
        text = re.sub(r"!\[.*?\]\(.*?\)", "", text)
        # Remove links (keep text)
        text = re.sub(r"\[([^\]]*)\]\(.*?\)", r"\1", text)
        # Remove headings markers
        text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
        # Remove bold/italic markers
        text = re.sub(r"[*_]{1,3}", "", text)
        # Remove horizontal rules
        text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)
        # Remove blockquotes
        text = re.sub(r"^>\s+", "", text, flags=re.MULTILINE)
        # Remove list markers
        text = re.sub(r"^[\s]*[-*+]\s+", "", text, flags=re.MULTILINE)
        text = re.sub(r"^\s*\d+\.\s+", "", text, flags=re.MULTILINE)
        # Normalize whitespace
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    @staticmethod
    def _extract_html(content: bytes) -> str:
        """Extract text from an HTML file.

        Strips HTML tags and returns visible text.
        """
        from html.parser import HTMLParser

        text = content.decode("utf-8", errors="replace")

        class _TextExtractHTMLParser(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self._text_parts: list[str] = []
                self._skip_tags = {"script", "style", "noscript", "svg"}

            def handle_data(self, data: str) -> None:
                stripped = data.strip()
                if stripped:
                    self._text_parts.append(stripped)

            def get_text(self) -> str:
                return " ".join(self._text_parts)

        parser = _TextExtractHTMLParser()
        parser.feed(text)
        return parser.get_text()

    @staticmethod
    def _extract_csv(content: bytes) -> str:
        """Extract text from a CSV file, returning a structured summary."""
        text = content.decode("utf-8", errors="replace")
        reader = csv.reader(StringIO(text))
        rows: list[str] = []
        for row in reader:
            rows.append(" | ".join(cell.strip() for cell in row))
        return "\n".join(rows)

    @staticmethod
    def _extract_json(content: bytes) -> str:
        """Extract text from a JSON file.

        Converts JSON to a readable string representation.
        """
        text = content.decode("utf-8", errors="replace")
        try:
            data = json.loads(text)
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            return text.strip()

    @staticmethod
    def _extract_xml(content: bytes) -> str:
        """Extract text from an XML file."""
        import xml.etree.ElementTree as ET

        text = content.decode("utf-8", errors="replace")
        try:
            root = ET.fromstring(text)

            def _iter_text(el: Any) -> list[str]:
                parts: list[str] = []
                if el.text and el.text.strip():
                    parts.append(el.text.strip())
                for child in el:
                    parts.extend(_iter_text(child))
                if el.tail and el.tail.strip():
                    parts.append(el.tail.strip())
                return parts

            return " ".join(_iter_text(root))
        except ET.ParseError:
            return text.strip()

    # ── Optional dependency extractors ───────────────────────────────

    async def _extract_optional(self, content: bytes, extension: str) -> str:
        """Extract text from formats that require optional dependencies.

        Args:
            content: Raw file bytes.
            extension: File extension (.pdf, .docx, .pptx).

        Returns:
            Extracted text.
        """
        if extension == ".pdf":
            return await self._extract_pdf(content)
        elif extension == ".docx":
            return self._extract_docx(content)
        elif extension == ".pptx":
            return self._extract_pptx(content)
        else:
            raise UnsupportedFileFormatError(extension)

    async def _extract_pdf(self, content: bytes) -> str:
        """Extract text from a PDF file using PyMuPDF or pdfminer.

        Tries PyMuPDF (fitz) first, then falls back to pdfminer.

        Args:
            content: Raw PDF bytes.

        Returns:
            Extracted text.
        """
        # Try PyMuPDF first
        try:
            import fitz  # type: ignore[import-untyped]

            doc = fitz.open(stream=content, filetype="pdf")
            pages: list[str] = []
            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)
                pages.append(f"[Page {page_num + 1}]\n{page.get_text()}")
            doc.close()
            return "\n\n".join(pages)

        except ImportError:
            logger.debug("PyMuPDF not available, trying pdfminer.six")

        # Fallback to pdfminer
        try:
            from io import BytesIO

            from pdfminer.high_level import extract_text_to_fp  # type: ignore[import-untyped]
            from pdfminer.layout import LAParams  # type: ignore[import-untyped]

            output = StringIO()
            with BytesIO(content) as file:
                extract_text_to_fp(file, output, laparams=LAParams(), output_type="text")
            return output.getvalue().strip()

        except ImportError:
            raise ImportError(
                "PDF text extraction requires either 'PyMuPDF' or 'pdfminer.six'. "
                "Install with: uv add PyMuPDF  # or: uv add pdfminer.six"
            )

    def _extract_docx(self, content: bytes) -> str:
        """Extract text from a .docx file using python-docx.

        Args:
            content: Raw .docx bytes.

        Returns:
            Extracted text.
        """
        try:
            from docx import Document  # type: ignore[import-untyped]
            from io import BytesIO

            doc = Document(BytesIO(content))
            paragraphs: list[str] = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        paragraphs.append(" | ".join(row_text))

            return "\n\n".join(paragraphs)

        except ImportError:
            raise ImportError(
                "DOCX text extraction requires 'python-docx'. "
                "Install with: uv add python-docx"
            )

    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from a .pptx file using python-pptx.

        Args:
            content: Raw .pptx bytes.

        Returns:
            Extracted text.
        """
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            from io import BytesIO

            prs = Presentation(BytesIO(content))
            slides: list[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            if any(row_text):
                                slide_texts.append(" | ".join(row_text))
                if slide_texts:
                    slides.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
            return "\n\n".join(slides)

        except ImportError:
            raise ImportError(
                "PPTX text extraction requires 'python-pptx'. "
                "Install with: uv add python-pptx"
            )

    # ── Utility ──────────────────────────────────────────────────────

    @staticmethod
    async def _read_file(path: Path) -> bytes:
        """Read a file asynchronously.

        Args:
            path: Path to the file.

        Returns:
            File content as bytes.
        """
        import asyncio

        return await asyncio.to_thread(path.read_bytes)

    @staticmethod
    def supports_extension(extension: str) -> bool:
        """Check if a file extension is supported.

        Args:
            extension: File extension (e.g., '.pdf').

        Returns:
            True if the format is supported.
        """
        return extension.lower() in SUPPORTED_EXTENSIONS or extension.lower() in OPTIONAL_EXTENSIONS
